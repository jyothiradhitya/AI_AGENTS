import pdfplumber
from pptx import Presentation
import docx
import pandas as pd
from pathlib import Path

def parse_pdf(path: str):
    texts = []
    with pdfplumber.open(path) as pdf:
        for i, page in enumerate(pdf.pages):
            t = page.extract_text() or ""
            texts.append({"page": i+1, "text": t})
    return texts

def parse_pptx(path: str):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
        slides.append({"slide": i+1, "text": "\n".join(text)})
    return slides

def parse_docx(path: str):
    doc = docx.Document(path)
    paragraphs = [p.text for p in doc.paragraphs]
    return [{"para_index": i, "text": p} for i, p in enumerate(paragraphs)]

def parse_csv(path: str):
    df = pd.read_csv(path, dtype=str, keep_default_na=False)
    chunks = []
    for i, row in df.iterrows():
        text = ", ".join([f"{c}:{row[c]}" for c in df.columns])
        chunks.append({"row": i, "text": text})
    return chunks

def parse_txt(path: str):
    txt = Path(path).read_text(encoding="utf-8")
    return [{"text": txt}]

def chunk_texts(text, chunk_size=500):
    chunks = []
    start = 0
    while start < len(text):
        chunk = text[start:start+chunk_size]
        chunks.append(chunk)
        start += chunk_size
    return chunks
