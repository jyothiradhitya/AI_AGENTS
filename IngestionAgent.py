import io
from datetime import datetime
from mcp_models import MCPMessage
import pdfplumber
from pptx import Presentation
import docx
import pandas as pd
from parsers import chunk_texts  # import your chunking function

class IngestionAgentSync:
    def __init__(self, out_queue=None):
        self.out_queue = out_queue
        self.name = "IngestionAgent"

    def log(self, text: str):
        print(f"[{self.name}] {text}")

    # ---- parsers ----
    def parse_pdf(self, uploaded_file):
        raw = uploaded_file.read()
        uploaded_file.seek(0)
        with pdfplumber.open(io.BytesIO(raw)) as pdf:
            text = "\n".join([p.extract_text() or "" for p in pdf.pages])
        return text

    def parse_pptx(self, uploaded_file):
        raw = uploaded_file.read()
        uploaded_file.seek(0)
        prs = Presentation(io.BytesIO(raw))
        text = "\n".join(
            "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            for slide in prs.slides
        )
        return text

    def parse_docx(self, uploaded_file):
        raw = uploaded_file.read()
        uploaded_file.seek(0)
        doc = docx.Document(io.BytesIO(raw))
        text = "\n".join([p.text for p in doc.paragraphs])
        return text

    def parse_csv(self, uploaded_file):
        raw = uploaded_file.read()
        uploaded_file.seek(0)
        df = pd.read_csv(io.BytesIO(raw), dtype=str, keep_default_na=False)
        text = "\n".join(
            [", ".join([f"{c}:{row[c]}" for c in df.columns]) for _, row in df.iterrows()]
        )
        return text

    def parse_txt(self, uploaded_file):
        raw = uploaded_file.read()
        uploaded_file.seek(0)
        text = raw.decode("utf-8", errors="ignore")
        return text

    # ---- sync handler ----
    def handle_sync(self, files, chunk_size=500):
        docs = []
        for f in files:
            name = getattr(f, "name", "unknown")
            if name.endswith(".pdf"):
                text = self.parse_pdf(f)
            elif name.endswith(".pptx"):
                text = self.parse_pptx(f)
            elif name.endswith(".docx"):
                text = self.parse_docx(f)
            elif name.endswith(".csv"):
                text = self.parse_csv(f)
            elif name.endswith((".txt", ".md")):
                text = self.parse_txt(f)
            else:
                raw = f.read()
                f.seek(0)
                try:
                    text = raw.decode("utf-8", errors="ignore")
                except Exception:
                    text = f"[Unsupported type {name}]"

            # ---- CHUNKING ----
            chunks = chunk_texts(text, chunk_size=chunk_size)
            docs.append({"filename": name, "chunks": chunks})

        return docs

    # ---- async handler ----
    async def handle(self, msg):
        files = msg.payload.get("files", [])
        docs = self.handle_sync(files)
        if self.out_queue is not None:
            ack = MCPMessage(
                trace_id=msg.trace_id,
                type="INGESTION_ACK",
                sender=self.name,
                receiver="CoordinatorAgent",
                timestamp=datetime.utcnow(),
                payload={"files": docs}
            )
            await self.out_queue.put(ack)
            self.log(f"Sent INGESTION_ACK with {len(docs)} chunked docs")
